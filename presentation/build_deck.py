"""
Build the Flixster presentation deck (.pptx) for a 5-minute class demo.

Setup (once):
    pip install python-pptx

Run:
    python build_deck.py
    # → writes Flixster_Presentation.pptx in this folder

Then: upload the .pptx to Google Drive → it opens as Google Slides.

Slides cover every required point from the assignment:
  what it does · features (incl. stretch) · enjoyed most ·
  most interesting/challenging · next steps / do differently · shout-outs.

⚠️ PERSONALIZE the lines tagged  # >>> YOURS  — the "enjoyed / challenging /
next steps / shout-outs" slides are your own words; the drafts below are a
strong starting point pulled from this project's real work, not gospel.
Also replace DEPLOY_URL once Render gives you the live link.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Brand palette (matches the app's locked Amazon Prime theme) ──────────────
BG        = RGBColor(0x0F, 0x17, 0x1E)   # --bg   page background
SURFACE   = RGBColor(0x1A, 0x24, 0x2F)   # --surface  card panels
SURFACE_2 = RGBColor(0x23, 0x2F, 0x3E)   # --surface-2
ACCENT    = RGBColor(0x00, 0xA8, 0xE1)   # --accent  Prime cyan
TEXT      = RGBColor(0xFF, 0xFF, 0xFF)   # --text
MUTED     = RGBColor(0xAA, 0xB7, 0xC4)   # --text-muted

# ── Personalize these ────────────────────────────────────────────────────────
PRESENTER  = "Ozias Tumimana"                      # >>> YOURS (pronouns?)
DEPLOY_URL = "https://<your-app>.onrender.com"     # >>> YOURS (Render link)
TAGLINE    = ("Flixster shows the movies playing in theaters right now — "
             "search them, save favorites, watch trailers, and get an AI "
             "recommendation, all in a streaming-service-style UI.")

# 16:9 widescreen
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ── Helpers ──────────────────────────────────────────────────────────────────
def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return s


def rect(s, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def textbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def set_run(r, text, size, color=TEXT, bold=False, italic=False):
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = "Arial"


def accent_bar(s, x=Inches(0.6), y=Inches(0.6), w=Inches(0.18), h=Inches(0.7)):
    """The small cyan tab that echoes the app's accent."""
    bar = rect(s, x, y, w, h, ACCENT)
    return bar


def heading(s, text, sub=None):
    accent_bar(s)
    tf = textbox(s, Inches(0.95), Inches(0.5), Inches(11.6), Inches(1.0))
    p = tf.paragraphs[0]
    set_run(p.add_run(), text, 34, TEXT, bold=True)
    if sub:
        tf2 = textbox(s, Inches(0.95), Inches(1.35), Inches(11.6), Inches(0.5))
        set_run(tf2.paragraphs[0].add_run(), sub, 16, MUTED)


def bullets(tf, items, size=18, gap=10):
    """items: list of (text, level) or plain strings (level 0)."""
    for i, item in enumerate(items):
        text, level = (item if isinstance(item, tuple) else (item, 0))
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        p.space_after = Pt(gap)
        marker = "✱  " if level == 0 else "–  "
        run = p.add_run()
        set_run(run, marker + text, size,
                TEXT if level == 0 else MUTED,
                bold=(level == 0))


def card(s, x, y, w, h, title, lines, title_color=ACCENT):
    rect(s, x, y, w, h, SURFACE)
    pad = Inches(0.3)
    tf = textbox(s, x + pad, y + pad, w - pad * 2, h - pad * 2)
    set_run(tf.paragraphs[0].add_run(), title, 20, title_color, bold=True)
    for line in lines:
        p = tf.add_paragraph()
        p.space_before = Pt(6)
        set_run(p.add_run(), line, 15, TEXT)


# ── Slide 1 — Title ──────────────────────────────────────────────────────────
s = slide()
# big cyan accent block on the left, like the hero's accent
rect(s, Inches(0), Inches(0), Inches(0.35), SH, ACCENT)
tf = textbox(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(2.0))
set_run(tf.paragraphs[0].add_run(), "Flixster", 72, TEXT, bold=True)
p = tf.add_paragraph()
set_run(p.add_run(), "Now playing in theaters — reimagined as a streaming app",
        22, ACCENT)
tf2 = textbox(s, Inches(0.9), Inches(4.4), Inches(11.5), Inches(1.5))
set_run(tf2.paragraphs[0].add_run(), f"by {PRESENTER}", 20, TEXT, bold=True)
p = tf2.add_paragraph()
set_run(p.add_run(), DEPLOY_URL, 18, MUTED)
p = tf2.add_paragraph()
p.space_before = Pt(4)
set_run(p.add_run(), "React 18 + Vite · TMDb API · OpenRouter AI · deployed on Render",
        15, MUTED, italic=True)


# ── Slide 2 — What it does ─────────────────────────────────────────────────────
s = slide()
heading(s, "What does Flixster do?")
tf = textbox(s, Inches(0.95), Inches(1.7), Inches(7.2), Inches(5.0))
set_run(tf.paragraphs[0].add_run(), TAGLINE, 20, TEXT)
bl = textbox(s, Inches(0.95), Inches(3.2), Inches(7.2), Inches(3.8))
bullets(bl, [
    "Pulls live “now playing” movies from The Movie Database (TMDb)",
    "Search any movie; sort by title, release date, or rating",
    "Click a movie for full details, genres, and an embedded trailer",
    "AI “Watch Recommendation” explains who would enjoy the film",
], size=18)
# side panel: the stack
card(s, Inches(8.5), Inches(1.7), Inches(4.0), Inches(5.0), "The stack", [
    "React 18 + Vite",
    "TMDb REST API",
    "OpenRouter (AI, free models)",
    "YouTube IFrame API (trailers)",
    "CSS custom-property theming",
    "Deployed: Render static site",
])


# ── Slide 3 — Features ─────────────────────────────────────────────────────────
s = slide()
heading(s, "Features", "Required spec + every stretch feature except none — plus extras")
card(s, Inches(0.95), Inches(1.9), Inches(3.7), Inches(4.9), "Required", [
    "Now-playing grid",
    "Search",
    "Load More pagination",
    "Sort (title / date / rating)",
    "Responsive 2/3/5-col layout",
    "Loading + error states",
], title_color=TEXT)
card(s, Inches(4.85), Inches(1.9), Inches(3.7), Inches(4.9), "Stretch", [
    "Favorites + Watched (Sets)",
    "Sidebar view-switcher",
    "Details modal",
    "Embedded YouTube trailers",
    "AI Watch Recommendation",
])
card(s, Inches(8.75), Inches(1.9), Inches(3.6), Inches(4.9), "Redesign extras", [
    "Auto-rotating hero billboard",
    "Category rows + genre chips",
    "Amazon Prime visual theme",
    "Mobile hamburger drawer",
    "Expanding search icon",
], title_color=MUTED)


# ── Slide 4 — Enjoyed most + Most challenging ─────────────────────────────────
s = slide()
heading(s, "Favorite & most challenging")
# Favorite
rect(s, Inches(0.95), Inches(1.8), Inches(5.5), Inches(5.0), SURFACE)
tf = textbox(s, Inches(1.25), Inches(2.1), Inches(4.9), Inches(4.4))
set_run(tf.paragraphs[0].add_run(), "💙  Enjoyed the most", 22, ACCENT, bold=True)
for line in [
    "Building the hero billboard — an auto-rotating, muted trailer behind "
    "the title, just like a real streaming home page.",
    "",
    "Seeing the whole app re-skin instantly the moment color tokens moved "
    "into CSS variables.",
    "",
    "# >>> YOURS: say which part was genuinely fun for you.",
]:
    p = tf.add_paragraph()
    p.space_before = Pt(8)
    set_run(p.add_run(), line, 16, TEXT if not line.startswith("#") else MUTED,
            italic=line.startswith("#"))
# Challenging
rect(s, Inches(6.85), Inches(1.8), Inches(5.5), Inches(5.0), SURFACE)
tf = textbox(s, Inches(7.15), Inches(2.1), Inches(4.9), Inches(4.4))
set_run(tf.paragraphs[0].add_run(), "🧩  Most interesting / challenging",
        22, ACCENT, bold=True)
for line in [
    "The YouTube trailer flashed its play/skip controls for a few seconds "
    "on load — fixed by hiding the video behind the poster until it’s "
    "actually playing, then revealing it.",
    "",
    "Keeping favorites as immutable Sets so React actually re-renders.",
    "",
    "# >>> YOURS: pick the one that taught you the most.",
]:
    p = tf.add_paragraph()
    p.space_before = Pt(8)
    set_run(p.add_run(), line, 16, TEXT if not line.startswith("#") else MUTED,
            italic=line.startswith("#"))


# ── Slide 5 — Next steps / do differently ─────────────────────────────────────
s = slide()
heading(s, "What I'd do next")
tf = textbox(s, Inches(0.95), Inches(1.9), Inches(11.4), Inches(4.8))
bullets(tf, [
    "Proxy the API keys through a small backend — right now VITE_ keys ship "
    "in the bundle (fine for free-tier, not for production)",
    "Persist favorites & watched to localStorage so they survive a reload",
    "Add a real router so movie details have shareable URLs",
    "Cache TMDb responses to cut repeat requests and speed up Load More",
    "# >>> YOURS: anything you'd architect differently a second time?",
], size=19, gap=14)


# ── Slide 6 — Shout-outs / resources ──────────────────────────────────────────
s = slide()
heading(s, "Thanks & resources")
tf = textbox(s, Inches(0.95), Inches(1.9), Inches(11.4), Inches(4.8))
bullets(tf, [
    "Shout-out to <person> for helping me with my project!   # >>> YOURS",
    "The Movie Database (TMDb) — movie data & images",
    "OpenRouter — free LLM access for the AI recommendation",
    "CodePath — course, spec, and feedback",
    "# >>> YOURS: add any tutorials, classmates, or docs you leaned on.",
], size=19, gap=14)
# footer link
foot = textbox(s, Inches(0.95), Inches(6.7), Inches(11.4), Inches(0.5))
set_run(foot.paragraphs[0].add_run(), DEPLOY_URL, 16, ACCENT)


prs.save("Flixster_Presentation.pptx")
print("✓ wrote Flixster_Presentation.pptx  (6 slides)")
